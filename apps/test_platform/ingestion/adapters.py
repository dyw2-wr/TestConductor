"""Deterministic adapters from uploaded bytes to requirement text.

Adapters use structured parsers for structured formats and never call a model,
execute embedded content, fetch links, or silently truncate extracted text.
"""

from __future__ import annotations

import csv
from io import BytesIO, StringIO
import json
import math
from pathlib import Path, PurePosixPath
import re
import shutil
from typing import Any, Mapping, Sequence
import zipfile

from apps.test_platform.intent.contracts import contains_secret_value

from .contracts import (
    AdapterResult,
    IngestionError,
    IngestionLimits,
    IngestionWarning,
    InputFile,
)


_SAFE_ID = re.compile(r"^[A-Za-z0-9][A-Za-z0-9_.@-]{0,127}$")
_TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".rst", ".org"}
_IMAGE_EXTENSIONS = {".bmp", ".jpeg", ".jpg", ".png", ".tif", ".tiff"}
_SUPPORTED_EXTENSIONS = (
    _TEXT_EXTENSIONS
    | _IMAGE_EXTENSIONS
    | {
        ".csv",
        ".docx",
        ".feature",
        ".htm",
        ".html",
        ".json",
        ".pdf",
        ".pptx",
        ".reqif",
        ".tsv",
        ".xls",
        ".xlsx",
        ".xml",
        ".yaml",
        ".yml",
    }
)
_LEGACY_OFFICE_EXTENSIONS = {".doc", ".ppt"}
_HTTP_METHODS = {
    "get",
    "put",
    "post",
    "delete",
    "options",
    "head",
    "patch",
    "trace",
}
_COMMON_HEADER_NAMES = {
    "acceptancecriteria",
    "content",
    "description",
    "issuekey",
    "key",
    "requirement",
    "requirementid",
    "requirements",
    "summary",
    "title",
    "workitemkey",
    "标题",
    "描述",
    "功能描述",
    "验收标准",
    "需求",
    "需求id",
    "需求描述",
    "需求编号",
}
_ID_HEADER_NAMES = {
    "id",
    "issuekey",
    "key",
    "requirementid",
    "workitemkey",
    "需求id",
    "需求编号",
}


def _json_text(value: Any, *, indent: int | None = None) -> str:
    """Serialize parser data without sorting mixed YAML key types."""

    try:
        return json.dumps(
            value,
            ensure_ascii=False,
            indent=indent,
            default=str,
            sort_keys=False,
        )
    except (TypeError, ValueError):
        return str(value)


def supported_extensions() -> tuple[str, ...]:
    return tuple(sorted(_SUPPORTED_EXTENSIONS))


def parse_input_file(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    """Select an adapter by extension and validate binary signatures where possible."""

    limits.validate()
    data = source.data
    if not data:
        raise IngestionError("EMPTY_FILE", "文件内容为空", source.filename)
    if len(data) > limits.max_file_bytes:
        raise IngestionError(
            "FILE_TOO_LARGE",
            f"文件超过 {limits.max_file_bytes} 字节限制",
            source.filename,
        )

    extension = Path(source.filename).suffix.lower()
    if extension in _LEGACY_OFFICE_EXTENSIONS:
        raise IngestionError(
            "LEGACY_OFFICE_UNSUPPORTED",
            f"旧格式 {extension} 需要先转换为 OOXML 格式",
            source.filename,
        )
    if extension not in _SUPPORTED_EXTENSIONS:
        raise IngestionError(
            "UNSUPPORTED_FORMAT",
            f"不支持的文件扩展名: {extension or '<none>'}",
            source.filename,
        )

    if extension in _TEXT_EXTENSIONS:
        result = _extract_plain_text(source, limits)
    elif extension == ".docx":
        result = _extract_docx(source, limits)
    elif extension == ".pdf":
        result = _extract_pdf(source, limits)
    elif extension in {".csv", ".tsv"}:
        result = _extract_delimited(source, limits, extension)
    elif extension == ".xlsx":
        result = _extract_xlsx(source, limits)
    elif extension == ".xls":
        result = _extract_xls(source, limits)
    elif extension in {".json", ".yaml", ".yml"}:
        result = _extract_json_or_yaml(source, limits, extension)
    elif extension == ".feature":
        result = _extract_gherkin(source, limits)
    elif extension in {".xml", ".reqif"}:
        result = _extract_xml(source, limits, force_reqif=extension == ".reqif")
    elif extension in {".html", ".htm"}:
        result = _extract_html(source, limits)
    elif extension == ".pptx":
        result = _extract_pptx(source, limits)
    else:
        result = _extract_image(source, limits)

    return _with_content_type_warning(result, source)


def _with_content_type_warning(result: AdapterResult, source: InputFile) -> AdapterResult:
    content_type = str(source.content_type or "").split(";", 1)[0].strip().lower()
    if not content_type or content_type == "application/octet-stream":
        return result
    expected = {
        "text": {"text/plain", "text/markdown", "text/x-rst"},
        "markdown": {"text/markdown", "text/plain"},
        "docx": {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        },
        "pdf": {"application/pdf"},
        "csv": {"text/csv", "text/plain", "application/csv"},
        "tsv": {"text/tab-separated-values", "text/plain"},
        "xlsx": {
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        },
        "xls": {"application/vnd.ms-excel"},
        "json": {"application/json", "text/json", "text/plain"},
        "yaml": {
            "application/yaml",
            "application/x-yaml",
            "text/yaml",
            "text/x-yaml",
            "text/plain",
        },
        "openapi": {
            "application/json",
            "application/yaml",
            "application/x-yaml",
            "text/yaml",
            "text/x-yaml",
            "text/plain",
        },
        "gherkin": {"text/plain", "text/x-gherkin"},
        "xml": {"application/xml", "text/xml", "text/plain"},
        "reqif": {"application/xml", "text/xml", "application/reqif+xml"},
        "html": {"text/html", "application/xhtml+xml"},
        "pptx": {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        },
        "image": {"image/bmp", "image/jpeg", "image/png", "image/tiff"},
    }.get(result.source_type, set())
    if expected and content_type not in expected:
        warning = IngestionWarning(
            code="CONTENT_TYPE_MISMATCH",
            message=(
                f"上传 content_type={content_type!r} 与扩展名解析结果"
                f" {result.source_type!r} 不一致；已按内容签名和扩展名处理"
            ),
            source_name=source.filename,
        )
        return AdapterResult(
            source_type=result.source_type,
            requirements=result.requirements,
            warnings=(*result.warnings, warning),
        )
    return result


def _decode_text(data: bytes, source_name: str) -> tuple[str, tuple[IngestionWarning, ...]]:
    warnings: list[IngestionWarning] = []
    if data.startswith(b"\xff\xfe") or data.startswith(b"\xfe\xff"):
        encoding = "utf-16"
        try:
            text = data.decode(encoding)
        except UnicodeDecodeError as exc:
            raise IngestionError(
                "TEXT_DECODE_FAILED",
                "UTF-16 文件解码失败",
                source_name,
            ) from exc
        warnings.append(
            IngestionWarning(
                "UTF16_BOM_DECODED",
                "文件使用 UTF-16 BOM 解码；请在预览中确认文字",
                source_name,
            )
        )
    else:
        if b"\x00" in data:
            raise IngestionError("NUL_BYTE_REJECTED", "文本包含 NUL 字节", source_name)
        try:
            text = data.decode("utf-8-sig")
            if data.startswith(b"\xef\xbb\xbf"):
                warnings.append(
                    IngestionWarning(
                        "UTF8_BOM_REMOVED",
                        "UTF-8 BOM 已在解码时移除",
                        source_name,
                    )
                )
        except UnicodeDecodeError:
            try:
                text = data.decode("gb18030")
            except UnicodeDecodeError as exc:
                raise IngestionError(
                    "TEXT_DECODE_FAILED",
                    "文件不是有效 UTF-8、UTF-16 或 GB18030 文本",
                    source_name,
                ) from exc
            warnings.append(
                IngestionWarning(
                    "ENCODING_FALLBACK",
                    "文件使用 GB18030 解码；请在预览中确认文字",
                    source_name,
                )
            )
    if any(
        (ord(character) < 32 and character not in "\t\r\n") or ord(character) == 127
        for character in text
    ):
        raise IngestionError(
            "BINARY_TEXT_REJECTED",
            "解码结果包含不可见控制字符，疑似二进制文件",
            source_name,
        )
    if not text.strip():
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "没有提取到非空文本", source_name)
    return text, tuple(warnings)


def _extract_plain_text(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    del limits
    text, warnings = _decode_text(source.data, source.filename)
    source_type = "markdown" if Path(source.filename).suffix.lower() in {".md", ".markdown"} else "text"
    return AdapterResult(source_type, ((None, text),), warnings)


def _validate_zip(
    source: InputFile,
    limits: IngestionLimits,
    *,
    required_member: str,
) -> tuple[zipfile.ZipFile, set[str], tuple[IngestionWarning, ...]]:
    if not source.data.startswith(b"PK"):
        raise IngestionError(
            "FILE_SIGNATURE_MISMATCH",
            "OOXML 文件缺少 ZIP 签名",
            source.filename,
        )
    try:
        archive = zipfile.ZipFile(BytesIO(source.data))
        infos = archive.infolist()
    except (OSError, zipfile.BadZipFile) as exc:
        raise IngestionError("INVALID_ZIP", "OOXML ZIP 容器损坏", source.filename) from exc
    try:
        if len(infos) > limits.max_archive_entries:
            raise IngestionError(
                "ARCHIVE_ENTRY_LIMIT",
                "OOXML ZIP 成员数量超过限制",
                source.filename,
            )
        names: set[str] = set()
        total_size = 0
        high_compression_entries: list[str] = []
        for info in infos:
            normalized = str(info.filename).replace("\\", "/")
            path = PurePosixPath(normalized)
            if path.is_absolute() or any(part in {"", ".", ".."} for part in path.parts):
                raise IngestionError(
                    "ARCHIVE_PATH_INVALID",
                    f"OOXML ZIP 包含越界成员: {info.filename}",
                    source.filename,
                )
            if normalized in names:
                raise IngestionError(
                    "ARCHIVE_DUPLICATE_ENTRY",
                    f"OOXML ZIP 包含重复成员: {normalized}",
                    source.filename,
                )
            names.add(normalized)
            if info.flag_bits & 0x1:
                raise IngestionError(
                    "ARCHIVE_ENCRYPTED",
                    "OOXML ZIP 包含加密成员",
                    source.filename,
                )
            if info.file_size > limits.max_archive_entry_bytes:
                raise IngestionError(
                    "ARCHIVE_ENTRY_TOO_LARGE",
                    f"OOXML ZIP 成员过大: {normalized}",
                    source.filename,
                )
            total_size += info.file_size
            if total_size > limits.max_archive_uncompressed_bytes:
                raise IngestionError(
                    "ARCHIVE_UNCOMPRESSED_LIMIT",
                    "OOXML ZIP 展开大小超过限制",
                    source.filename,
                )
            if normalized.lower().endswith((".xml", ".rels")):
                member_bytes = archive.read(info)
                if re.search(
                    br"<!\s*(?:DOCTYPE|ENTITY)\b",
                    member_bytes,
                    flags=re.IGNORECASE,
                ):
                    raise IngestionError(
                        "XML_DTD_REJECTED",
                        "OOXML XML 成员包含不允许的 DTD/ENTITY",
                        source.filename,
                    )
            if info.file_size and (
                info.compress_size == 0
                or info.file_size / max(1, info.compress_size) > limits.max_archive_ratio
            ):
                # Repetitive requirement prose compresses very well.  Keep the
                # hard byte/member limits above, and only reject an extreme
                # ratio that is also large enough to resemble a zip bomb.
                ratio = info.file_size / max(1, info.compress_size)
                if info.file_size >= 1024 * 1024 and ratio > limits.max_archive_ratio * 100:
                    raise IngestionError(
                        "ARCHIVE_COMPRESSION_RATIO",
                        f"OOXML ZIP 压缩比异常: {normalized}",
                        source.filename,
                    )
                high_compression_entries.append(normalized)
        if required_member not in names:
            raise IngestionError(
                "OOXML_PART_MISSING",
                f"OOXML 缺少必要成员: {required_member}",
                source.filename,
            )
        lowered = {name.lower() for name in names}
        if any("vbaproject.bin" in name for name in lowered):
            raise IngestionError(
                "ACTIVE_CONTENT_REJECTED",
                "包含 VBA 宏的 Office 文件不能作为需求输入",
                source.filename,
            )
        if any("/embeddings/" in f"/{name}" for name in lowered):
            raise IngestionError(
                "EMBEDDED_OBJECT_REJECTED",
                "包含嵌入对象的 Office 文件不能作为需求输入",
                source.filename,
            )
        warnings: list[IngestionWarning] = []
        if high_compression_entries:
            warnings.append(
                IngestionWarning(
                    "ARCHIVE_COMPRESSION_RATIO_HIGH",
                    "OOXML 包含高压缩比正文；已通过展开大小门禁，请在预览中确认内容",
                    source.filename,
                )
            )
        external_relation = False
        for info in infos:
            if not info.filename.lower().endswith(".rels"):
                continue
            relation = archive.read(info)
            if b'TargetMode="External"' in relation or b"TargetMode='External'" in relation:
                external_relation = True
                break
        if external_relation:
            warnings.append(
                IngestionWarning(
                    "EXTERNAL_LINKS_IGNORED",
                    "Office 外部链接目标已忽略，不会发起网络请求",
                    source.filename,
                )
            )
        return archive, names, tuple(warnings)
    except Exception:
        archive.close()
        raise


def _extract_docx(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    archive, names, archive_warnings = _validate_zip(
        source, limits, required_member="word/document.xml"
    )
    document_xml = archive.read("word/document.xml")
    archive.close()
    try:
        from docx import Document
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        document = Document(BytesIO(source.data))
        lines: list[str] = []

        def append_blocks(parent: Any) -> None:
            for child in parent.iterchildren():
                if isinstance(child, CT_P):
                    text = Paragraph(child, document).text
                    if text.strip():
                        lines.append(text)
                elif isinstance(child, CT_Tbl):
                    table = Table(child, document)
                    for row in table.rows:
                        values = [
                            cell.text.replace("\n", " / ").strip()
                            for cell in row.cells
                        ]
                        if any(values):
                            lines.append(" | ".join(values))
                else:
                    # Content controls and similar wrappers may contain normal
                    # paragraphs/tables.  Recurse only through the wrapper;
                    # table/paragraph descendants are handled exactly once.
                    append_blocks(child)

        append_blocks(document.element.body)

        for textbox in document.element.body.xpath(".//w:txbxContent"):
            textbox_text = "".join(
                str(node.text or "") for node in textbox.xpath(".//w:t")
            ).strip()
            if textbox_text and textbox_text not in lines:
                lines.append(f"[Text box] {textbox_text}")

        seen_parts: set[str] = set()
        header_footer_included = False
        for section in document.sections:
            for label, container in (
                ("Header", section.header),
                ("Footer", section.footer),
            ):
                part_name = str(container.part.partname)
                if part_name in seen_parts:
                    continue
                seen_parts.add(part_name)
                part_lines = [
                    paragraph.text
                    for paragraph in container.paragraphs
                    if paragraph.text.strip()
                ]
                for table in container.tables:
                    for row in table.rows:
                        values = [
                            cell.text.replace("\n", " / ").strip()
                            for cell in row.cells
                        ]
                        if any(values):
                            part_lines.append(" | ".join(values))
                if part_lines:
                    header_footer_included = True
                    lines.append(f"[{label}] " + "\n".join(part_lines))
    except Exception as exc:
        raise IngestionError("DOCX_PARSE_FAILED", "DOCX 解析失败", source.filename) from exc
    text = "\n".join(lines)
    if not text.strip():
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "DOCX 没有可用正文", source.filename)
    warnings = list(archive_warnings)
    lowered = {name.lower() for name in names}
    if any(name.startswith("word/media/") for name in lowered):
        warnings.append(
            IngestionWarning(
                "EMBEDDED_IMAGES_IGNORED",
                "DOCX 图片未做 OCR；仅提取正文和表格文字",
                source.filename,
            )
        )
    ignored_parts = {
        part
        for part in ("word/comments.xml", "word/footnotes.xml", "word/endnotes.xml")
        if part in lowered
    }
    if ignored_parts:
        warnings.append(
            IngestionWarning(
                "AUXILIARY_PARTS_IGNORED",
                "DOCX 批注、脚注或尾注未进入需求正文",
                source.filename,
            )
        )
    numbering_detected = b"<w:numPr" in document_xml or any(
        "numPr" in str(paragraph.style.element.xml)
        for paragraph in document.paragraphs
    )
    if numbering_detected:
        warnings.append(
            IngestionWarning(
                "WORD_NUMBERING_LABELS_NOT_RENDERED",
                "DOCX 自动编号的正文已保留，但 Word 渲染出的序号可能未进入纯文本",
                source.filename,
            )
        )
    if header_footer_included:
        warnings.append(
            IngestionWarning(
                "HEADER_FOOTER_INCLUDED",
                "DOCX 页眉/页脚文字已带标签加入正文；请确认它是否属于测试需求",
                source.filename,
            )
        )
    return AdapterResult("docx", ((None, text),), tuple(warnings))


def _extract_pdf(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    if b"%PDF-" not in source.data[:1024]:
        raise IngestionError(
            "FILE_SIGNATURE_MISMATCH",
            "PDF 文件缺少 %PDF 签名",
            source.filename,
        )
    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(source.data), strict=False)
        if reader.is_encrypted and reader.decrypt("") == 0:
            raise IngestionError(
                "PDF_ENCRYPTED",
                "PDF 已加密，需要先解除密码保护",
                source.filename,
            )
        if len(reader.pages) > limits.max_pdf_pages:
            raise IngestionError(
                "PDF_PAGE_LIMIT",
                f"PDF 页数超过 {limits.max_pdf_pages} 页限制",
                source.filename,
            )
        page_texts: list[str] = []
        empty_pages = 0
        for page in reader.pages:
            value = page.extract_text() or ""
            page_texts.append(value)
            if not value.strip():
                empty_pages += 1
    except IngestionError:
        raise
    except Exception as exc:
        raise IngestionError("PDF_PARSE_FAILED", "PDF 文本提取失败", source.filename) from exc

    warnings: list[IngestionWarning] = []
    text = "\n\n".join(value for value in page_texts if value.strip())
    if not text.strip():
        text, ocr_warning = _ocr_pdf(source, limits)
        warnings.append(ocr_warning)
    elif empty_pages:
        warnings.append(
            IngestionWarning(
                "PDF_PAGES_WITHOUT_TEXT",
                f"PDF 有 {empty_pages} 页未提取到文字；请检查扫描页或图片内容",
                source.filename,
            )
        )
    return AdapterResult("pdf", ((None, text),), tuple(warnings))


def _ocr_pdf(
    source: InputFile, limits: IngestionLimits
) -> tuple[str, IngestionWarning]:
    if shutil.which("tesseract") is None:
        raise IngestionError(
            "OCR_UNAVAILABLE",
            "PDF 没有文本层，当前运行环境未安装 Tesseract OCR",
            source.filename,
        )
    try:
        from pdf2image import convert_from_bytes
        from pypdf import PdfReader
        import pytesseract

        reader = PdfReader(BytesIO(source.data), strict=False)
        pages = len(reader.pages)
        if pages < 1 or pages > limits.max_ocr_pages:
            raise IngestionError(
                "OCR_PAGE_LIMIT",
                f"扫描 PDF OCR 最多处理 {limits.max_ocr_pages} 页",
                source.filename,
            )
        dpi = 200
        estimated_pixels = 0
        for page in reader.pages:
            width = max(1.0, float(page.mediabox.width))
            height = max(1.0, float(page.mediabox.height))
            estimated_pixels += int(width / 72 * dpi) * int(height / 72 * dpi)
            if estimated_pixels > limits.max_ocr_pixels:
                raise IngestionError(
                    "OCR_PIXEL_LIMIT",
                    "扫描 PDF OCR 像素总量超过限制",
                    source.filename,
                )
        languages = set(pytesseract.get_languages(config=""))
        preferred = [value for value in ("chi_sim", "eng") if value in languages]
        kwargs = {"lang": "+".join(preferred)} if preferred else {}
        page_texts: list[str] = []
        actual_pixels = 0
        for page_number in range(1, pages + 1):
            images = convert_from_bytes(
                source.data,
                dpi=dpi,
                first_page=page_number,
                last_page=page_number,
            )
            try:
                for image in images:
                    actual_pixels += image.width * image.height
                    if actual_pixels > limits.max_ocr_pixels:
                        raise IngestionError(
                            "OCR_PIXEL_LIMIT",
                            "扫描 PDF OCR 像素总量超过限制",
                            source.filename,
                        )
                    page_texts.append(pytesseract.image_to_string(image, **kwargs))
            finally:
                for image in images:
                    image.close()
        text = "\n\n".join(page_texts)
    except IngestionError:
        raise
    except Exception as exc:
        raise IngestionError("OCR_FAILED", "扫描 PDF OCR 失败", source.filename) from exc
    if not text.strip():
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "OCR 没有识别出文字", source.filename)
    return text, IngestionWarning(
        "OCR_USED",
        "PDF 没有文本层，已使用 OCR；必须人工核对识别结果",
        source.filename,
    )


def _extract_image(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    if shutil.which("tesseract") is None:
        raise IngestionError(
            "OCR_UNAVAILABLE",
            "图片需求需要 Tesseract OCR，当前运行环境未安装",
            source.filename,
        )
    try:
        from PIL import Image, ImageSequence
        import pytesseract

        image = Image.open(BytesIO(source.data))
        image.verify()
        image = Image.open(BytesIO(source.data))
        languages = set(pytesseract.get_languages(config=""))
        preferred = [value for value in ("chi_sim", "eng") if value in languages]
        kwargs = {"lang": "+".join(preferred)} if preferred else {}
        texts: list[str] = []
        total_pixels = 0
        frame_count = 0
        for frame in ImageSequence.Iterator(image):
            frame_count += 1
            if frame_count > limits.max_ocr_pages:
                raise IngestionError(
                    "OCR_PAGE_LIMIT",
                    f"多帧图片 OCR 最多处理 {limits.max_ocr_pages} 帧",
                    source.filename,
                )
            frame_pixels = frame.width * frame.height
            total_pixels += frame_pixels
            if total_pixels > limits.max_ocr_pixels:
                raise IngestionError(
                    "OCR_PIXEL_LIMIT",
                    "图片像素总量超过 OCR 限制",
                    source.filename,
                )
            frame = frame.convert("RGB")
            texts.append(pytesseract.image_to_string(frame, **kwargs))
            frame.close()
        image.close()
        text = "\n\n".join(texts)
    except IngestionError:
        raise
    except Exception as exc:
        raise IngestionError("IMAGE_PARSE_FAILED", "图片 OCR 失败", source.filename) from exc
    if not text.strip():
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "图片 OCR 没有识别出文字", source.filename)
    warning = IngestionWarning(
        "OCR_USED",
        "图片已使用 OCR；必须人工核对识别结果",
        source.filename,
    )
    return AdapterResult("image", ((None, text),), (warning,))


def _extract_html(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    del limits
    raw, decode_warnings = _decode_text(source.data, source.filename)
    try:
        from bs4 import BeautifulSoup, Comment

        soup = BeautifulSoup(raw, "html.parser")
        for element in soup(["script", "style", "noscript", "template"]):
            element.decompose()
        for comment in soup.find_all(string=lambda item: isinstance(item, Comment)):
            comment.extract()
        lines = [line.strip() for line in soup.get_text("\n").splitlines() if line.strip()]
        text = "\n".join(lines)
    except Exception as exc:
        raise IngestionError("HTML_PARSE_FAILED", "HTML 解析失败", source.filename) from exc
    if not text.strip():
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "HTML 没有可见文字", source.filename)
    return AdapterResult("html", ((None, text),), decode_warnings)


def _extract_pptx(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    archive, names, archive_warnings = _validate_zip(
        source, limits, required_member="ppt/presentation.xml"
    )
    archive.close()
    try:
        from pptx import Presentation

        presentation = Presentation(BytesIO(source.data))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = str(getattr(shape, "text", "") or "")
                if text.strip():
                    lines.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [cell.text.replace("\n", " / ").strip() for cell in row.cells]
                        if any(values):
                            lines.append(" | ".join(values))
    except Exception as exc:
        raise IngestionError("PPTX_PARSE_FAILED", "PPTX 解析失败", source.filename) from exc
    text = "\n".join(lines)
    if not text.strip():
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "PPTX 没有可用文字", source.filename)
    warnings = list(archive_warnings)
    if any(name.lower().startswith("ppt/media/") for name in names):
        warnings.append(
            IngestionWarning(
                "EMBEDDED_IMAGES_IGNORED",
                "PPTX 图片未做 OCR；仅提取文本和表格",
                source.filename,
            )
        )
    return AdapterResult("pptx", ((None, text),), tuple(warnings))


def _normalize_header(value: Any) -> str:
    return re.sub(r"[\s_.:/\\-]+", "", str(value or "").strip().lower())


def _cell_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "true" if value else "false"
    if isinstance(value, float) and math.isfinite(value) and value.is_integer():
        return str(int(value))
    return str(value)


def _rows_to_requirements(
    rows: Sequence[Sequence[Any]],
    *,
    source_name: str,
    source_label: str,
    max_requirements: int | None = None,
) -> tuple[list[tuple[str | None, str]], list[IngestionWarning]]:
    normalized_rows = [
        [_cell_text(value) for value in row]
        for row in rows
        if any(_cell_text(value).strip() for value in row)
    ]
    if not normalized_rows:
        return [], []
    width = max(len(row) for row in normalized_rows)
    normalized_rows = [row + [""] * (width - len(row)) for row in normalized_rows]
    candidate_headers = normalized_rows[0]
    normalized_headers = [_normalize_header(value) for value in candidate_headers]
    has_header = any(value in _COMMON_HEADER_NAMES for value in normalized_headers)
    warnings: list[IngestionWarning] = []
    if has_header:
        headers = [value.strip() or f"column_{index}" for index, value in enumerate(candidate_headers, start=1)]
        data_rows = normalized_rows[1:]
    else:
        headers = [f"column_{index}" for index in range(1, width + 1)]
        data_rows = normalized_rows
        warnings.append(
            IngestionWarning(
                "TABLE_HEADER_NOT_DETECTED",
                f"{source_label} 未识别出常见表头，已按 column_1... 保留每行",
                source_name,
            )
        )
    requirements: list[tuple[str | None, str]] = []
    for row in data_rows:
        pairs = [(headers[index], value) for index, value in enumerate(row) if value.strip()]
        if not pairs:
            continue
        content = "\n".join(f"{header}: {value}" for header, value in pairs)
        explicit_id = None
        for index, header in enumerate(normalized_headers if has_header else []):
            if header in _ID_HEADER_NAMES and index < len(row):
                candidate = row[index].strip()
                if _SAFE_ID.fullmatch(candidate):
                    explicit_id = candidate
                    break
        requirements.append((explicit_id, content))
        if max_requirements is not None and len(requirements) > max_requirements:
            raise IngestionError(
                "TOO_MANY_REQUIREMENTS",
                f"{source_label} 提取出的需求超过 {max_requirements} 条",
                source_name,
            )
    return requirements, warnings


def _extract_delimited(
    source: InputFile, limits: IngestionLimits, extension: str
) -> AdapterResult:
    text, decode_warnings = _decode_text(source.data, source.filename)
    try:
        # The stdlib default is 128 KiB, below the v4 per-requirement limit.
        # Raise it monotonically so a legitimate long cell is not truncated or
        # rejected before the service applies its own byte limits.
        csv.field_size_limit(max(csv.field_size_limit(), limits.max_requirement_bytes))
        if extension == ".tsv":
            dialect = csv.excel_tab
        else:
            try:
                dialect = csv.Sniffer().sniff(text[:8192], delimiters=",;\t|")
            except csv.Error:
                dialect = csv.excel
        rows = []
        for row in csv.reader(StringIO(text), dialect=dialect):
            if not any(str(value).strip() for value in row):
                continue
            rows.append(row)
            # One row may be a header.  At max+2 non-empty rows the file
            # necessarily contains more requirements than v4 can accept.
            if len(rows) > limits.max_requirements + 1:
                raise IngestionError(
                    "TOO_MANY_REQUIREMENTS",
                    f"CSV/TSV 提取出的需求超过 {limits.max_requirements} 条",
                    source.filename,
                )
    except (csv.Error, UnicodeError) as exc:
        raise IngestionError("CSV_PARSE_FAILED", "CSV/TSV 解析失败", source.filename) from exc
    cell_count = sum(len(row) for row in rows)
    if cell_count > limits.max_table_cells:
        raise IngestionError(
            "TABLE_CELL_LIMIT",
            f"表格单元格超过 {limits.max_table_cells} 限制",
            source.filename,
        )
    requirements, row_warnings = _rows_to_requirements(
        rows,
        source_name=source.filename,
        source_label="CSV/TSV",
        max_requirements=limits.max_requirements,
    )
    if not requirements:
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "CSV/TSV 没有数据行", source.filename)
    if any(
        str(value).lstrip().startswith("=")
        for row in rows
        for value in row
    ):
        row_warnings.append(
            IngestionWarning(
                "FORMULAS_NOT_EXECUTED",
                "CSV/TSV 中以 '=' 开头的内容只作为文字保留",
                source.filename,
            )
        )
    source_type = "tsv" if extension == ".tsv" else "csv"
    return AdapterResult(
        source_type,
        tuple(requirements),
        (*decode_warnings, *row_warnings),
    )


def _extract_xlsx(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    archive, names, archive_warnings = _validate_zip(
        source, limits, required_member="xl/workbook.xml"
    )
    archive.close()
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(
            BytesIO(source.data),
            read_only=True,
            data_only=False,
            keep_links=False,
        )
        requirements: list[tuple[str | None, str]] = []
        warnings: list[IngestionWarning] = list(archive_warnings)
        cell_count = 0
        formula_found = False
        for sheet in workbook.worksheets:
            if sheet.sheet_state != "visible":
                warnings.append(
                    IngestionWarning(
                        "HIDDEN_SHEET_IGNORED",
                        f"隐藏工作表 {sheet.title!r} 未进入需求正文",
                        source.filename,
                    )
                )
                continue
            rows: list[list[Any]] = []
            remaining = limits.max_requirements - len(requirements)
            for row in sheet.iter_rows(values_only=True):
                values = list(row)
                cell_count += len(values)
                if cell_count > limits.max_table_cells:
                    raise IngestionError(
                        "TABLE_CELL_LIMIT",
                        f"工作簿单元格超过 {limits.max_table_cells} 限制",
                        source.filename,
                    )
                formula_found = formula_found or any(
                    isinstance(value, str) and value.startswith("=") for value in values
                )
                if not any(_cell_text(value).strip() for value in values):
                    continue
                rows.append(values)
                if len(rows) > remaining + 1:
                    raise IngestionError(
                        "TOO_MANY_REQUIREMENTS",
                        f"工作簿提取出的需求超过 {limits.max_requirements} 条",
                        source.filename,
                    )
            extracted, row_warnings = _rows_to_requirements(
                rows,
                source_name=source.filename,
                source_label=f"工作表 {sheet.title!r}",
                max_requirements=max(0, limits.max_requirements - len(requirements)),
            )
            requirements.extend(extracted)
            warnings.extend(row_warnings)
        workbook.close()
    except IngestionError:
        raise
    except Exception as exc:
        raise IngestionError("XLSX_PARSE_FAILED", "XLSX 解析失败", source.filename) from exc
    if not requirements:
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "XLSX 没有可用数据行", source.filename)
    if formula_found:
        warnings.append(
            IngestionWarning(
                "FORMULAS_NOT_EXECUTED",
                "Excel 公式只作为文字保留，适配器不会计算公式",
                source.filename,
            )
        )
    if any(name.lower().startswith("xl/externallinks/") for name in names):
        warnings.append(
            IngestionWarning(
                "EXTERNAL_LINKS_IGNORED",
                "Excel 外部链接已忽略，不会读取外部工作簿",
                source.filename,
            )
        )
    return AdapterResult("xlsx", tuple(requirements), tuple(warnings))


def _extract_xls(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    if not source.data.startswith(bytes.fromhex("D0CF11E0A1B11AE1")):
        raise IngestionError(
            "FILE_SIGNATURE_MISMATCH",
            "XLS 文件缺少 OLE2 签名",
            source.filename,
        )
    try:
        import olefile

        with olefile.OleFileIO(BytesIO(source.data)) as ole:
            stream_names = {
                "/".join(parts).lower() for parts in ole.listdir(streams=True)
            }
        if any("vba" in name for name in stream_names):
            raise IngestionError(
                "ACTIVE_CONTENT_REJECTED",
                "包含 VBA 宏的 XLS 文件不能作为需求输入",
                source.filename,
            )
        if any(name.startswith("objectpool/") for name in stream_names):
            raise IngestionError(
                "EMBEDDED_OBJECT_REJECTED",
                "包含嵌入对象的 XLS 文件不能作为需求输入",
                source.filename,
            )

        import xlrd

        workbook = xlrd.open_workbook(file_contents=source.data, on_demand=True)
        requirements: list[tuple[str | None, str]] = []
        warnings: list[IngestionWarning] = []
        cell_count = 0
        for sheet in workbook.sheets():
            if getattr(sheet, "visibility", 0) != 0:
                warnings.append(
                    IngestionWarning(
                        "HIDDEN_SHEET_IGNORED",
                        f"隐藏工作表 {sheet.name!r} 未进入需求正文",
                        source.filename,
                    )
                )
                continue
            rows = []
            remaining = limits.max_requirements - len(requirements)
            for row_index in range(sheet.nrows):
                row = sheet.row_values(row_index)
                cell_count += len(row)
                if cell_count > limits.max_table_cells:
                    raise IngestionError(
                        "TABLE_CELL_LIMIT",
                        f"工作簿单元格超过 {limits.max_table_cells} 限制",
                        source.filename,
                    )
                if not any(_cell_text(value).strip() for value in row):
                    continue
                rows.append(row)
                if len(rows) > remaining + 1:
                    raise IngestionError(
                        "TOO_MANY_REQUIREMENTS",
                        f"工作簿提取出的需求超过 {limits.max_requirements} 条",
                        source.filename,
                    )
            extracted, row_warnings = _rows_to_requirements(
                rows,
                source_name=source.filename,
                source_label=f"工作表 {sheet.name!r}",
                max_requirements=max(0, limits.max_requirements - len(requirements)),
            )
            requirements.extend(extracted)
            warnings.extend(row_warnings)
        workbook.release_resources()
    except IngestionError:
        raise
    except Exception as exc:
        raise IngestionError("XLS_PARSE_FAILED", "XLS 解析失败", source.filename) from exc
    if not requirements:
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "XLS 没有可用数据行", source.filename)
    warnings.append(
        IngestionWarning(
            "LEGACY_EXCEL_VALUES_ONLY",
            "XLS 只提取工作簿保存的单元格值；不会执行宏或重新计算公式",
            source.filename,
        )
    )
    return AdapterResult("xls", tuple(requirements), tuple(warnings))


def _check_structure(value: Any, limits: IngestionLimits, source_name: str) -> None:
    stack: list[tuple[Any, int]] = [(value, 1)]
    nodes = 0
    while stack:
        current, depth = stack.pop()
        nodes += 1
        if nodes > limits.max_json_nodes:
            raise IngestionError(
                "STRUCTURE_NODE_LIMIT",
                "JSON/YAML 节点数量超过限制",
                source_name,
            )
        if depth > limits.max_json_depth:
            raise IngestionError(
                "STRUCTURE_DEPTH_LIMIT",
                "JSON/YAML 嵌套深度超过限制",
                source_name,
            )
        if isinstance(current, Mapping):
            stack.extend((item, depth + 1) for item in current.values())
        elif isinstance(current, (list, tuple)):
            stack.extend((item, depth + 1) for item in current)


def _extract_json_or_yaml(
    source: InputFile, limits: IngestionLimits, extension: str
) -> AdapterResult:
    text, decode_warnings = _decode_text(source.data, source.filename)
    try:
        if extension == ".json":
            value = json.loads(text)
            source_type = "json"
        else:
            import yaml

            value = yaml.safe_load(text)
            source_type = "yaml"
    except Exception as exc:
        raise IngestionError(
            "STRUCTURED_TEXT_PARSE_FAILED",
            f"{extension} 解析失败",
            source.filename,
        ) from exc
    _check_structure(value, limits, source.filename)
    if contains_secret_value(value):
        raise IngestionError(
            "SECRET_LITERAL_REJECTED",
            "JSON/YAML 包含疑似密码、token、私钥或 API key 实际值",
            source.filename,
        )
    if isinstance(value, Mapping) and (
        isinstance(value.get("openapi"), str) or isinstance(value.get("swagger"), str)
    ):
        result = _extract_openapi(
            value, source.filename, max_requirements=limits.max_requirements
        )
        return AdapterResult(
            result.source_type,
            result.requirements,
            (*decode_warnings, *result.warnings),
        )
    if isinstance(value, Mapping) and isinstance(value.get("requirements"), list):
        requirements, record_warnings = _extract_requirement_records(
            value["requirements"],
            source.filename,
            max_requirements=limits.max_requirements,
        )
        return AdapterResult(
            source_type,
            tuple(requirements),
            (*decode_warnings, *record_warnings),
        )
    if isinstance(value, list):
        requirements, record_warnings = _extract_requirement_records(
            value,
            source.filename,
            max_requirements=limits.max_requirements,
        )
        warning = IngestionWarning(
            "GENERIC_RECORDS_SPLIT",
            "顶层数组已按记录拆分；请在预览中确认每项是否为独立需求",
            source.filename,
        )
        return AdapterResult(
            source_type,
            tuple(requirements),
            (*decode_warnings, warning, *record_warnings),
        )
    warning = IngestionWarning(
        "STRUCTURED_TEXT_PASSED_THROUGH",
        "未识别专用需求 schema，已把完整 JSON/YAML 原文作为一条需求",
        source.filename,
    )
    return AdapterResult(source_type, ((None, text),), (*decode_warnings, warning))


def _extract_requirement_records(
    values: Sequence[Any], source_name: str, max_requirements: int | None = None
) -> tuple[list[tuple[str | None, str]], list[IngestionWarning]]:
    requirements: list[tuple[str | None, str]] = []
    warnings: list[IngestionWarning] = []
    for value in values:
        explicit_id = None
        if isinstance(value, str):
            content = value
        elif isinstance(value, Mapping):
            candidate_id = value.get("requirement_id") or value.get("id") or value.get("key")
            if candidate_id is not None and _SAFE_ID.fullmatch(str(candidate_id).strip()):
                explicit_id = str(candidate_id).strip()
            if isinstance(value.get("content"), str) and set(value).issubset(
                {"requirement_id", "content"}
            ):
                content = value["content"]
            else:
                content = _json_text(value, indent=2)
        else:
            content = _json_text(value)
        if not str(content).strip():
            warnings.append(
                IngestionWarning(
                    "EMPTY_RECORD_IGNORED",
                    "结构化文件中的空记录已忽略",
                    source_name,
                )
            )
            continue
        requirements.append((explicit_id, str(content)))
        if max_requirements is not None and len(requirements) > max_requirements:
            raise IngestionError(
                "TOO_MANY_REQUIREMENTS",
                f"结构化文件提取出的需求超过 {max_requirements} 条",
                source_name,
            )
    if not requirements:
        raise IngestionError(
            "EMPTY_EXTRACTED_TEXT",
            "结构化文件没有可用需求记录",
            source_name,
        )
    return requirements, warnings


def _extract_openapi(
    value: Mapping[str, Any], source_name: str, max_requirements: int | None = None
) -> AdapterResult:
    paths = value.get("paths")
    if not isinstance(paths, Mapping):
        raise IngestionError("OPENAPI_PATHS_MISSING", "OpenAPI 缺少 paths 对象", source_name)
    requirements: list[tuple[str | None, str]] = []
    unresolved_refs = False
    shared_local_refs: dict[str, Any] = {}
    version = value.get("openapi") or value.get("swagger")
    global_security = value.get("security")
    global_servers = value.get("servers") or value.get("host")
    for path, path_item in paths.items():
        if not isinstance(path_item, Mapping):
            continue
        path_parameters = path_item.get("parameters", [])
        for method, operation in path_item.items():
            method_name = str(method).lower()
            if method_name not in _HTTP_METHODS or not isinstance(operation, Mapping):
                continue
            lines = [f"OpenAPI {version}", f"{method_name.upper()} {path}"]
            for key, label in (
                ("summary", "Summary"),
                ("description", "Description"),
                ("deprecated", "Deprecated"),
            ):
                if key in operation and operation[key] is not None and operation[key] != "":
                    value_text = (
                        operation[key]
                        if isinstance(operation[key], str)
                        else _json_text(operation[key])
                    )
                    lines.append(f"{label}: {value_text}")
            parameters = [
                *list(path_parameters if isinstance(path_parameters, list) else []),
                *list(operation.get("parameters", []) if isinstance(operation.get("parameters"), list) else []),
            ]
            if parameters:
                lines.append(
                    "Parameters:\n"
                    + _json_text(parameters, indent=2)
                )
            if operation.get("requestBody") is not None:
                lines.append(
                    "Request body:\n"
                    + _json_text(operation.get("requestBody"), indent=2)
                )
            responses = operation.get("responses")
            if responses is not None:
                lines.append(
                    "Responses:\n"
                    + _json_text(responses, indent=2)
                )
            security = operation.get("security", global_security)
            if security is not None:
                lines.append(
                    "Security:\n"
                    + _json_text(security, indent=2)
                )
            servers = operation.get("servers", global_servers)
            if servers is not None:
                lines.append(
                    "Servers:\n"
                    + _json_text(servers, indent=2)
                )
            local_refs, has_unresolved_ref = _openapi_reference_context(
                value,
                {
                    "path_parameters": path_parameters,
                    "operation": operation,
                },
                source_name,
            )
            if local_refs:
                shared_local_refs.update(local_refs)
                lines.append("Local refs used: " + ", ".join(local_refs))
            operation_id = operation.get("operationId")
            explicit_id = (
                str(operation_id).strip()
                if operation_id is not None
                and _SAFE_ID.fullmatch(str(operation_id).strip())
                else None
            )
            unresolved_refs = unresolved_refs or has_unresolved_ref
            requirements.append((explicit_id, "\n".join(lines)))
            if max_requirements is not None and len(requirements) > max_requirements:
                raise IngestionError(
                    "TOO_MANY_REQUIREMENTS",
                    f"OpenAPI operation 超过 {max_requirements} 条",
                    source_name,
                )
    if not requirements:
        raise IngestionError(
            "OPENAPI_NO_OPERATIONS",
            "OpenAPI paths 中没有可用 HTTP operation",
            source_name,
        )
    if shared_local_refs:
        definitions = ["OpenAPI local referenced definitions (shared by this file):"]
        for ref, resolved in shared_local_refs.items():
            definitions.append(f"{ref}:\n{_json_text(resolved, indent=2)}")
        first_id, first_content = requirements[0]
        requirements[0] = (first_id, "\n".join(definitions) + "\n\n" + first_content)
    warnings = (
        (
            IngestionWarning(
                "OPENAPI_REFS_NOT_RESOLVED",
                "OpenAPI $ref 只作为引用文字保留，不读取外部文件或 URL",
                source_name,
            ),
        )
        if unresolved_refs
        else ()
    )
    return AdapterResult("openapi", tuple(requirements), warnings)


_MISSING_REFERENCE = object()


def _find_openapi_refs(value: Any) -> set[str]:
    refs: set[str] = set()
    stack = [value]
    while stack:
        current = stack.pop()
        if isinstance(current, Mapping):
            for key, item in current.items():
                if key == "$ref" and isinstance(item, str):
                    refs.add(item)
                stack.append(item)
        elif isinstance(current, (list, tuple)):
            stack.extend(current)
    return refs


def _resolve_openapi_pointer(root: Any, ref: str) -> Any:
    if not ref.startswith("#/"):
        return _MISSING_REFERENCE
    current = root
    for token in ref[2:].split("/"):
        token = token.replace("~1", "/").replace("~0", "~")
        if isinstance(current, Mapping) and token in current:
            current = current[token]
        elif isinstance(current, list) and token.isdigit() and int(token) < len(current):
            current = current[int(token)]
        else:
            return _MISSING_REFERENCE
    return current


def _openapi_reference_context(
    root: Mapping[str, Any],
    operation_context: Mapping[str, Any],
    source_name: str,
) -> tuple[dict[str, Any], bool]:
    """Resolve bounded local refs while leaving external refs as warnings."""

    pending = sorted(_find_openapi_refs(operation_context))
    resolved: dict[str, Any] = {}
    unresolved = False
    examined: set[str] = set()
    while pending:
        ref = pending.pop(0)
        if ref in examined:
            continue
        examined.add(ref)
        if len(examined) > 100:
            raise IngestionError(
                "OPENAPI_REF_LIMIT",
                "单个 OpenAPI operation 的本地引用超过 100 个",
                source_name,
            )
        value = _resolve_openapi_pointer(root, ref)
        if value is _MISSING_REFERENCE:
            unresolved = True
            continue
        resolved[ref] = value
        for nested in sorted(_find_openapi_refs(value)):
            if nested not in examined:
                pending.append(nested)
    return resolved, unresolved


def _extract_gherkin(source: InputFile, limits: IngestionLimits) -> AdapterResult:
    text, decode_warnings = _decode_text(source.data, source.filename)
    try:
        from gherkin.parser import Parser

        document = Parser().parse(text)
    except ImportError as exc:
        raise IngestionError(
            "GHERKIN_PARSER_UNAVAILABLE",
            "缺少 gherkin-official 解析器",
            source.filename,
        ) from exc
    except Exception as exc:
        raise IngestionError("GHERKIN_PARSE_FAILED", "Gherkin 解析失败", source.filename) from exc
    feature = document.get("feature") if isinstance(document, Mapping) else None
    if not isinstance(feature, Mapping):
        raise IngestionError("GHERKIN_FEATURE_MISSING", "Gherkin 缺少 Feature", source.filename)

    feature_background: list[Mapping[str, Any]] = []
    scenarios: list[tuple[Mapping[str, Any], Mapping[str, Any] | None, list[Mapping[str, Any]]]] = []
    for child in feature.get("children", []) or []:
        if not isinstance(child, Mapping):
            continue
        if isinstance(child.get("background"), Mapping):
            feature_background = list(child["background"].get("steps", []) or [])
        elif isinstance(child.get("scenario"), Mapping):
            scenarios.append((child["scenario"], None, feature_background))
            if len(scenarios) > limits.max_requirements:
                raise IngestionError(
                    "TOO_MANY_REQUIREMENTS",
                    f"Gherkin scenario 超过 {limits.max_requirements} 条",
                    source.filename,
                )
        elif isinstance(child.get("rule"), Mapping):
            rule = child["rule"]
            rule_background: list[Mapping[str, Any]] = []
            for rule_child in rule.get("children", []) or []:
                if not isinstance(rule_child, Mapping):
                    continue
                if isinstance(rule_child.get("background"), Mapping):
                    rule_background = list(rule_child["background"].get("steps", []) or [])
                elif isinstance(rule_child.get("scenario"), Mapping):
                    scenarios.append(
                        (
                            rule_child["scenario"],
                            rule,
                            [*feature_background, *rule_background],
                        )
                    )
                    if len(scenarios) > limits.max_requirements:
                        raise IngestionError(
                            "TOO_MANY_REQUIREMENTS",
                            f"Gherkin scenario 超过 {limits.max_requirements} 条",
                            source.filename,
                        )
    if not scenarios:
        raise IngestionError(
            "GHERKIN_SCENARIO_MISSING",
            "Gherkin Feature 中没有 Scenario/Example",
            source.filename,
        )

    requirements: list[tuple[str | None, str]] = []
    feature_tags = _gherkin_tag_names(feature)
    for scenario, rule, background_steps in scenarios:
        lines = [f"{feature.get('keyword', 'Feature')}: {feature.get('name', '')}"]
        if feature_tags:
            lines.append("Feature tags: " + " ".join(feature_tags))
        description = str(feature.get("description") or "").strip()
        if description:
            lines.append(description)
        if rule is not None:
            rule_tags = _gherkin_tag_names(rule)
            if rule_tags:
                lines.append("Rule tags: " + " ".join(rule_tags))
            lines.append(f"{rule.get('keyword', 'Rule')}: {rule.get('name', '')}")
            rule_description = str(rule.get("description") or "").strip()
            if rule_description:
                lines.append(rule_description)
        tags = _gherkin_tag_names(scenario)
        if tags:
            lines.append("Tags: " + " ".join(tag for tag in tags if tag))
        lines.append(f"{scenario.get('keyword', 'Scenario')}: {scenario.get('name', '')}")
        scenario_description = str(scenario.get("description") or "").strip()
        if scenario_description:
            lines.append(scenario_description)
        if background_steps:
            lines.append("Background:")
            lines.extend(_render_gherkin_step(step) for step in background_steps)
        lines.extend(
            _render_gherkin_step(step) for step in scenario.get("steps", []) or []
        )
        for example in scenario.get("examples", []) or []:
            example_tags = _gherkin_tag_names(example)
            if example_tags:
                lines.append("Examples tags: " + " ".join(example_tags))
            lines.append(f"{example.get('keyword', 'Examples')}: {example.get('name', '')}")
            header = [cell.get("value", "") for cell in (example.get("tableHeader") or {}).get("cells", [])]
            if header:
                lines.append(" | ".join(str(value) for value in header))
            for row in example.get("tableBody", []) or []:
                values = [cell.get("value", "") for cell in row.get("cells", [])]
                lines.append(" | ".join(str(value) for value in values))
        requirements.append((None, "\n".join(line for line in lines if line is not None)))
    return AdapterResult("gherkin", tuple(requirements), decode_warnings)


def _gherkin_tag_names(node: Mapping[str, Any]) -> list[str]:
    return [
        str(item.get("name") or "")
        for item in node.get("tags", []) or []
        if isinstance(item, Mapping) and str(item.get("name") or "")
    ]


def _render_gherkin_step(step: Mapping[str, Any]) -> str:
    line = f"{step.get('keyword', '')}{step.get('text', '')}".rstrip()
    data_table = step.get("dataTable")
    if isinstance(data_table, Mapping):
        rows = []
        for row in data_table.get("rows", []) or []:
            values = [cell.get("value", "") for cell in row.get("cells", [])]
            rows.append("  | " + " | ".join(str(value) for value in values) + " |")
        if rows:
            line += "\n" + "\n".join(rows)
    doc_string = step.get("docString")
    if isinstance(doc_string, Mapping):
        media_type = str(doc_string.get("mediaType") or "")
        content = str(doc_string.get("content") or "")
        line += f"\n  ```{media_type}\n{content}\n  ```"
    return line


def _extract_xml(
    source: InputFile,
    limits: IngestionLimits,
    *,
    force_reqif: bool,
) -> AdapterResult:
    if re.search(br"<!\s*(?:DOCTYPE|ENTITY)\b", source.data, flags=re.IGNORECASE):
        raise IngestionError(
            "XML_DTD_REJECTED",
            "XML DTD/ENTITY 不允许进入需求摄取",
            source.filename,
        )
    try:
        from lxml import etree

        parser = etree.XMLParser(
            resolve_entities=False,
            no_network=True,
            load_dtd=False,
            huge_tree=False,
            recover=False,
        )
        root = etree.fromstring(source.data, parser=parser)
        docinfo = root.getroottree().docinfo
        if docinfo.doctype or docinfo.internalDTD is not None or docinfo.externalDTD is not None:
            raise IngestionError(
                "XML_DTD_REJECTED",
                "XML DTD/ENTITY 不允许进入需求摄取",
                source.filename,
            )
    except IngestionError:
        raise
    except Exception as exc:
        raise IngestionError("XML_PARSE_FAILED", "XML/ReqIF 解析失败", source.filename) from exc
    nodes = 0
    stack = [(root, 1)]
    while stack:
        element, depth = stack.pop()
        nodes += 1
        if nodes > limits.max_json_nodes:
            raise IngestionError("STRUCTURE_NODE_LIMIT", "XML 节点数量超过限制", source.filename)
        if depth > limits.max_json_depth:
            raise IngestionError("STRUCTURE_DEPTH_LIMIT", "XML 嵌套深度超过限制", source.filename)
        stack.extend((child, depth + 1) for child in element)
    local_name = _xml_local_name(root.tag)
    if force_reqif or local_name == "REQ-IF":
        return _extract_reqif(
            root,
            source.filename,
            max_requirements=limits.max_requirements,
        )
    lines = [value.strip() for value in root.itertext() if value and value.strip()]
    text = "\n".join(lines)
    if not text:
        raise IngestionError("EMPTY_EXTRACTED_TEXT", "XML 没有文本内容", source.filename)
    warning = IngestionWarning(
        "GENERIC_XML_TEXT",
        "未识别 ReqIF，已按安全 XML 可见文本提取",
        source.filename,
    )
    return AdapterResult("xml", ((None, text),), (warning,))


def _xml_local_name(tag: Any) -> str:
    return str(tag).rsplit("}", 1)[-1]


def _extract_reqif(
    root: Any, source_name: str, max_requirements: int | None = None
) -> AdapterResult:
    definitions: dict[str, str] = {}
    for element in root.iter():
        if not _xml_local_name(element.tag).startswith("ATTRIBUTE-DEFINITION-"):
            continue
        identifier = str(element.get("IDENTIFIER") or "")
        if identifier:
            definitions[identifier] = str(element.get("LONG-NAME") or identifier)

    requirements: list[tuple[str | None, str]] = []
    for spec_object in root.iter():
        if _xml_local_name(spec_object.tag) != "SPEC-OBJECT":
            continue
        identifier = str(spec_object.get("IDENTIFIER") or "").strip()
        lines: list[str] = []
        long_name = str(spec_object.get("LONG-NAME") or "").strip()
        if long_name:
            lines.append(f"Title: {long_name}")
        for value in spec_object.iter():
            value_type = _xml_local_name(value.tag)
            if not value_type.startswith("ATTRIBUTE-VALUE-"):
                continue
            definition_ref = None
            for child in value.iter():
                if _xml_local_name(child.tag).endswith("-REF") and child.text:
                    definition_ref = child.text.strip()
                    break
            label = definitions.get(definition_ref or "", definition_ref or value_type)
            extracted = str(value.get("THE-VALUE") or "").strip()
            if not extracted:
                for child in value:
                    if _xml_local_name(child.tag) == "THE-VALUE":
                        extracted = " ".join(
                            piece.strip()
                            for piece in child.itertext()
                            if piece and piece.strip()
                        )
                        break
            if not extracted:
                refs = [
                    child.text.strip()
                    for child in value.iter()
                    if _xml_local_name(child.tag).endswith("-REF")
                    and child.text
                    and child.text.strip() != definition_ref
                ]
                extracted = ", ".join(refs)
            if extracted:
                lines.append(f"{label}: {extracted}")
        if not lines:
            continue
        explicit_id = identifier if _SAFE_ID.fullmatch(identifier) else None
        requirements.append((explicit_id, "\n".join(lines)))
        if max_requirements is not None and len(requirements) > max_requirements:
            raise IngestionError(
                "TOO_MANY_REQUIREMENTS",
                f"ReqIF SPEC-OBJECT 超过 {max_requirements} 条",
                source_name,
            )
    if not requirements:
        raise IngestionError(
            "REQIF_NO_REQUIREMENTS",
            "ReqIF 中没有可用 SPEC-OBJECT",
            source_name,
        )
    warnings: list[IngestionWarning] = []
    if any(_xml_local_name(element.tag) == "SPEC-RELATION" for element in root.iter()):
        warnings.append(
            IngestionWarning(
                "REQIF_RELATIONS_NOT_PROJECTED",
                "ReqIF 关系未进入 v4 RequirementInput；正文与对象 ID 已保留",
                source_name,
            )
        )
    return AdapterResult("reqif", tuple(requirements), tuple(warnings))


__all__ = ["parse_input_file", "supported_extensions"]
