from __future__ import annotations

from io import BytesIO
import json
from types import SimpleNamespace
import unittest
from unittest.mock import patch
import zipfile

from apps.test_platform.ingestion import (
    IngestionError,
    IngestionLimits,
    InputFile,
    prepare_request,
)
from apps.test_platform.intent.builder import DefaultDesignBuilder
from apps.test_platform.intent.contracts import (
    DesignSelections,
    TargetSelection,
    contains_secret_literal,
    contains_secret_value,
)
from apps.test_platform.intent.prompt_builder import DefaultDesignPromptBuilder
from apps.test_platform.intent.service import TestDesignPipeline
from apps.test_platform.workflow import IntentToExecutionWorkflow


def _target() -> TargetSelection:
    return TargetSelection(system_id="ingestion-demo", environment="test")


def _selections() -> DesignSelections:
    return DesignSelections(techniques=["positive"], allowed_channels=["api"])


def _docx_bytes(*, external_link: bool = False) -> bytes:
    from docx import Document

    document = Document()
    document.add_paragraph("Health endpoint requirements")
    if external_link:
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.opc.constants import RELATIONSHIP_TYPE as RT

        paragraph = document.add_paragraph()
        relationship_id = paragraph.part.relate_to(
            "https://example.invalid/never-fetch",
            RT.HYPERLINK,
            is_external=True,
        )
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), relationship_id)
        run = OxmlElement("w:r")
        text = OxmlElement("w:t")
        text.text = "reviewed link label"
        run.append(text)
        hyperlink.append(run)
        paragraph._p.append(hyperlink)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Rule"
    table.cell(0, 1).text = "Expected"
    table.cell(1, 0).text = "GET /health"
    table.cell(1, 1).text = "200"
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _rich_docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.sections[0].header.paragraphs[0].text = "Account service PRD"
    document.add_paragraph("A" * 100_000)
    document.add_paragraph("Login requirement", style="List Number")
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Health requirements"
    table = slide.shapes.add_table(2, 2, 0, 0, 4_000_000, 1_000_000).table
    table.cell(0, 0).text = "Rule"
    table.cell(0, 1).text = "Expected"
    table.cell(1, 0).text = "GET /health"
    table.cell(1, 1).text = "200"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _pdf_bytes(text: str | None = None) -> bytes:
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    page = writer.add_blank_page(width=300, height=300)
    if text is not None:
        font = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Font"),
                NameObject("/Subtype"): NameObject("/Type1"),
                NameObject("/BaseFont"): NameObject("/Helvetica"),
            }
        )
        font_ref = writer._add_object(font)
        page[NameObject("/Resources")] = DictionaryObject(
            {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})}
        )
        stream = DecodedStreamObject()
        escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        stream.set_data(
            f"BT /F1 12 Tf 40 250 Td ({escaped}) Tj ET".encode("ascii")
        )
        page[NameObject("/Contents")] = writer._add_object(stream)
    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Requirements"
    sheet.append(["Requirement ID", "Summary", "Description"])
    sheet.append(["REQ-API-1", "Health", "Returns 200"])
    sheet.append(["REQ-API-2", "Latency", "p95 <= 500ms"])
    sheet.append([None, "Formula", "=1+1"])
    hidden = workbook.create_sheet("Hidden")
    hidden.sheet_state = "hidden"
    hidden.append(["Summary", "Description"])
    hidden.append(["Hidden", "must not appear"])
    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def _xls_bytes() -> bytes:
    import xlwt

    workbook = xlwt.Workbook()
    sheet = workbook.add_sheet("Requirements")
    values = [
        ["Requirement ID", "Summary", "Description"],
        ["REQ-XLS-1", "Database", "status is ready"],
    ]
    for row_index, row in enumerate(values):
        for column_index, value in enumerate(row):
            sheet.write(row_index, column_index, value)
    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


class _FixtureGateway:
    def generate(self, messages, output_schema):
        return {
            "title": "Ingested requirement",
            "objective": {"text": "Validate the ingested requirement"},
            "in_scope": [{"text": "Health endpoint"}],
            "out_of_scope": [],
            "scenarios": [
                {
                    "title": "Health succeeds",
                    "techniques": ["positive"],
                    "requirement_ids": ["REQ-0001"],
                    "required_states": [],
                    "operations": [
                        {"text": "Request the health endpoint", "channel_hint": "api"}
                    ],
                    "expected_results": [
                        {
                            "text": "The endpoint returns success",
                            "after_operation_index": 1,
                            "channel_hint": "api",
                        }
                    ],
                    "data_requirements": [],
                    "state_impact": {
                        "impact": "read_only",
                        "rationale": {"text": "The endpoint is read only"},
                    },
                }
            ],
            "open_questions": [],
        }


class RequirementIngestionV4Tests(unittest.TestCase):
    def test_result_schema_and_invalid_caller_values_use_stable_errors(self):
        result = prepare_request(
            frontend_text="health returns 200",
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(result.as_dict()["schema_version"], "requirement-ingestion.v1")

        with self.assertRaises(IngestionError) as invalid_target:
            prepare_request(
                frontend_text="health returns 200",
                target={"system_id": "", "environment": "test"},
                selections=_selections(),
            )
        self.assertEqual(invalid_target.exception.code, "REQUEST_INVALID")

        with self.assertRaises(IngestionError) as invalid_file:
            prepare_request(
                files=[{"filename": "bad.txt", "data": "not bytes"}],
                target=_target(),
                selections=_selections(),
            )
        self.assertEqual(invalid_file.exception.code, "FILE_INPUT_INVALID")

    def test_utf16_and_frontend_unicode_boundaries_are_explicit(self):
        result = prepare_request(
            files=[InputFile("requirement.txt", "健康接口返回 200".encode("utf-16"))],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(result.request.requirements[0].content, "健康接口返回 200")
        self.assertIn("UTF16_BOM_DECODED", [item.code for item in result.warnings])

        for raw, expected_code in (
            ("bad\ud800unicode", "INVALID_UNICODE"),
            ("clear\x1b[2J", "CONTROL_CHARACTER_REJECTED"),
        ):
            with self.subTest(expected_code=expected_code):
                with self.assertRaises(IngestionError) as raised:
                    prepare_request(
                        frontend_text=raw,
                        target=_target(),
                        selections=_selections(),
                    )
                self.assertEqual(raised.exception.code, expected_code)

    def test_frontend_text_is_exact_and_mixed_input_order_is_stable(self):
        frontend_text = "  用户直接输入\r\n健康接口返回 200  "
        result = prepare_request(
            frontend_text=frontend_text,
            files=[InputFile("notes.md", "# 补充\n数据库状态为 ready".encode("utf-8"))],
            target=_target(),
            selections=_selections(),
        )

        self.assertEqual(len(result.request.requirements), 2)
        self.assertEqual(result.request.requirements[0].content, frontend_text)
        self.assertEqual(result.request.requirements[1].content, "# 补充\n数据库状态为 ready")
        self.assertEqual([item.source_type for item in result.sources], ["frontend_text", "markdown"])
        self.assertIsNone(result.request.requirements[0].requirement_id)

    def test_frontend_text_reaches_first_layer_without_trimming(self):
        raw = "  no heading\r\njust one frontend paragraph  "
        prepared = prepare_request(
            frontend_text=raw,
            target=_target(),
            selections=_selections(),
        )
        pipeline = TestDesignPipeline(
            DefaultDesignBuilder(DefaultDesignPromptBuilder(), _FixtureGateway())
        )

        generated = pipeline.generate(prepared.request)

        self.assertEqual(generated.request.requirements[0].content, raw)
        self.assertEqual(generated.request.requirements[0].requirement_id, "REQ-0001")
        self.assertEqual(generated.input_snapshot.requirements[0].content, raw)

    def test_empty_frontend_text_is_ignored_only_when_a_file_exists(self):
        result = prepare_request(
            frontend_text="   \r\n",
            files=[InputFile("requirement.txt", b"health returns 200")],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(len(result.request.requirements), 1)
        self.assertIn("EMPTY_TEXT_IGNORED", [item.code for item in result.warnings])

        with self.assertRaisesRegex(IngestionError, "NO_INPUT"):
            prepare_request(
                frontend_text="   ",
                target=_target(),
                selections=_selections(),
            )

    def test_docx_extracts_paragraphs_tables_and_ignores_external_target(self):
        result = prepare_request(
            files=[InputFile("requirements.docx", _docx_bytes(external_link=True))],
            target=_target(),
            selections=_selections(),
        )
        content = result.request.requirements[0].content
        self.assertIn("Health endpoint requirements", content)
        self.assertIn("reviewed link label", content)
        self.assertIn("Rule | Expected", content)
        self.assertNotIn("example.invalid", content)
        self.assertIn("EXTERNAL_LINKS_IGNORED", [item.code for item in result.warnings])

    def test_docx_high_compression_header_and_numbering_are_not_silent(self):
        result = prepare_request(
            files=[InputFile("large-requirement.docx", _rich_docx_bytes())],
            target=_target(),
            selections=_selections(),
        )
        content = result.request.requirements[0].content
        codes = [item.code for item in result.warnings]
        self.assertIn("A" * 1000, content)
        self.assertIn("[Header] Account service PRD", content)
        self.assertIn("ARCHIVE_COMPRESSION_RATIO_HIGH", codes)
        self.assertIn("HEADER_FOOTER_INCLUDED", codes)
        self.assertIn("WORD_NUMBERING_LABELS_NOT_RENDERED", codes)

    def test_pdf_text_is_extracted_and_image_only_pdf_requires_ocr(self):
        result = prepare_request(
            files=[InputFile("requirements.pdf", _pdf_bytes("Health returns 200"))],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(result.request.requirements[0].content.strip(), "Health returns 200")

        with patch("apps.test_platform.ingestion.adapters.shutil.which", return_value=None):
            with self.assertRaises(IngestionError) as raised:
                prepare_request(
                    files=[InputFile("scan.pdf", _pdf_bytes())],
                    target=_target(),
                    selections=_selections(),
                )
        self.assertEqual(raised.exception.code, "OCR_UNAVAILABLE")

    def test_image_input_reports_missing_ocr_engine_before_parsing(self):
        with patch("apps.test_platform.ingestion.adapters.shutil.which", return_value=None):
            with self.assertRaises(IngestionError) as raised:
                prepare_request(
                    files=[InputFile("requirement.png", b"not-decoded-without-ocr")],
                    target=_target(),
                    selections=_selections(),
                )
        self.assertEqual(raised.exception.code, "OCR_UNAVAILABLE")
        self.assertEqual(raised.exception.source_name, "requirement.png")

    def test_csv_xlsx_and_xls_are_split_by_explicit_rows(self):
        csv_result = prepare_request(
            files=[
                InputFile(
                    "jira.csv",
                    (
                        "Requirement ID,Summary,Description\n"
                        "REQ-CSV-1,Health,Returns 200\n"
                        "REQ-CSV-2,Database,Status is ready\n"
                    ).encode("utf-8"),
                )
            ],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(
            [item.requirement_id for item in csv_result.request.requirements],
            ["REQ-CSV-1", "REQ-CSV-2"],
        )
        self.assertIn("Summary: Health", csv_result.request.requirements[0].content)

        xlsx_result = prepare_request(
            files=[InputFile("requirements.xlsx", _xlsx_bytes())],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(len(xlsx_result.request.requirements), 3)
        codes = [item.code for item in xlsx_result.warnings]
        self.assertIn("HIDDEN_SHEET_IGNORED", codes)
        self.assertIn("FORMULAS_NOT_EXECUTED", codes)
        self.assertNotIn("must not appear", json.dumps(xlsx_result.as_dict()))

        xls_result = prepare_request(
            files=[InputFile("requirements.xls", _xls_bytes())],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(xls_result.request.requirements[0].requirement_id, "REQ-XLS-1")
        self.assertIn("status is ready", xls_result.request.requirements[0].content)

    def test_long_csv_field_and_pptx_are_supported_without_truncation(self):
        long_value = "x" * 150_000
        csv_result = prepare_request(
            files=[
                InputFile(
                    "long.csv",
                    f"Summary,Description\nLong,{long_value}\n".encode("utf-8"),
                )
            ],
            target=_target(),
            selections=_selections(),
        )
        self.assertIn(long_value, csv_result.request.requirements[0].content)

        pptx_result = prepare_request(
            files=[InputFile("requirements.pptx", _pptx_bytes())],
            target=_target(),
            selections=_selections(),
        )
        self.assertIn("Health requirements", pptx_result.request.requirements[0].content)
        self.assertIn("GET /health | 200", pptx_result.request.requirements[0].content)

    def test_openapi_json_and_yaml_are_split_by_operation(self):
        document = {
            "openapi": "3.1.0",
            "paths": {
                "/health": {
                    "get": {
                        "operationId": "getHealth",
                        "summary": "Read health",
                        "responses": {
                            "200": {
                                "description": "ready",
                                "content": {
                                    "application/json": {
                                        "schema": {"$ref": "#/components/schemas/Health"}
                                    }
                                },
                            }
                        },
                    }
                },
                "/accounts": {
                    "post": {
                        "operationId": "createAccount",
                        "requestBody": {"required": True},
                        "responses": {"201": {"description": "created"}},
                    }
                },
            },
        }
        result = prepare_request(
            files=[InputFile("openapi.json", json.dumps(document).encode("utf-8"))],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(
            [item.requirement_id for item in result.request.requirements],
            ["getHealth", "createAccount"],
        )
        self.assertIn("GET /health", result.request.requirements[0].content)
        self.assertIn("POST /accounts", result.request.requirements[1].content)
        self.assertIn("OPENAPI_REFS_NOT_RESOLVED", [item.code for item in result.warnings])

        yaml_text = """openapi: 3.0.0
paths:
  /ready:
    get:
      responses:
        '200':
          description: ready
"""
        yaml_result = prepare_request(
            files=[InputFile("openapi.yaml", yaml_text.encode("utf-8"))],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(yaml_result.sources[0].source_type, "openapi")
        self.assertIn("GET /ready", yaml_result.request.requirements[0].content)

    def test_openapi_shared_local_schema_is_included_once_and_secrets_are_distinguished(self):
        paths = {
            f"/accounts/{index}": {
                "post": {
                    "operationId": f"createAccount{index}",
                    "summary": [] if index == 0 else f"Create account {index}",
                    "requestBody": {
                        "content": {
                            "application/json": {
                                "schema": {"$ref": "#/components/schemas/Login"}
                            }
                        }
                    },
                    "responses": {"201": {"description": "created"}},
                }
            }
            for index in range(20)
        }
        document = {
            "openapi": "3.1.0",
            "paths": paths,
            "components": {
                "schemas": {
                    "Login": {
                        "type": "object",
                        "properties": {
                            "password": {
                                "type": "string",
                                "format": "password",
                                "minLength": 8,
                            },
                            "padding": {"type": "string", "description": "x" * 15_000},
                        },
                    }
                }
            },
        }
        result = prepare_request(
            files=[InputFile("login-openapi.json", json.dumps(document).encode("utf-8"))],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(len(result.request.requirements), 20)
        combined = "\n".join(item.content for item in result.request.requirements)
        self.assertEqual(combined.count("OpenAPI local referenced definitions"), 1)
        self.assertIn('"minLength": 8', combined)
        self.assertNotIn("OPENAPI_REFS_NOT_RESOLVED", [item.code for item in result.warnings])
        self.assertFalse(contains_secret_literal('"password": {"type": "string"}'))
        self.assertFalse(
            contains_secret_value({"password": {"type": "string", "format": "password"}})
        )
        self.assertFalse(
            contains_secret_value(
                {"name": "password", "source": "secret", "secret": True}
            )
        )
        self.assertFalse(contains_secret_literal('{"secret": false}'))
        self.assertTrue(contains_secret_literal('{"password": "hunter2"}'))

        pipeline = TestDesignPipeline(
            DefaultDesignBuilder(DefaultDesignPromptBuilder(), _FixtureGateway())
        )
        pipeline.generate(result.request)

        document["components"]["schemas"]["Login"]["properties"]["password"][
            "example"
        ] = "hunter2"
        with self.assertRaises(IngestionError) as secret_error:
            prepare_request(
                files=[
                    InputFile(
                        "unsafe-openapi.json",
                        json.dumps(document).encode("utf-8"),
                    )
                ],
                target=_target(),
                selections=_selections(),
            )
        self.assertEqual(secret_error.exception.code, "SECRET_LITERAL_REJECTED")

    def test_gherkin_uses_official_parser_and_preserves_scenario_structure(self):
        feature = """Feature: Health
  Background:
    Given the service is running

  Scenario: Ready
    When GET /health
    Then status is 200

  Scenario Outline: Invalid path
    When GET <path>
    Then status is <status>
    Examples:
      | path     | status |
      | /missing | 404    |
"""
        result = prepare_request(
            files=[InputFile("health.feature", feature.encode("utf-8"))],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(len(result.request.requirements), 2)
        self.assertIn("Background:", result.request.requirements[0].content)
        self.assertIn("Scenario: Ready", result.request.requirements[0].content)
        self.assertIn("Examples:", result.request.requirements[1].content)
        self.assertIn("/missing | 404", result.request.requirements[1].content)

    def test_gherkin_preserves_feature_rule_scenario_and_examples_tags(self):
        feature = """@feature_tag
Feature: Accounts
  @rule_tag
  Rule: Password policy
    @scenario_tag
    Scenario Outline: Validate password
      When password is <value>
      Then result is <result>
      @examples_tag
      Examples: values
        | value | result |
        | short | reject |
"""
        result = prepare_request(
            files=[InputFile("accounts.feature", feature.encode("utf-8"))],
            target=_target(),
            selections=_selections(),
        )
        content = result.request.requirements[0].content
        self.assertIn("Feature tags: @feature_tag", content)
        self.assertIn("Rule tags: @rule_tag", content)
        self.assertIn("Tags: @scenario_tag", content)
        self.assertIn("Examples tags: @examples_tag", content)

    def test_reqif_objects_are_split_and_generic_xml_rejects_dtd(self):
        reqif = b'''<?xml version="1.0"?>
<REQ-IF>
  <CORE-CONTENT><REQ-IF-CONTENT>
    <SPEC-TYPES><SPEC-OBJECT-TYPE IDENTIFIER="TYPE-1"><SPEC-ATTRIBUTES>
      <ATTRIBUTE-DEFINITION-STRING IDENTIFIER="ATTR-1" LONG-NAME="Description"/>
    </SPEC-ATTRIBUTES></SPEC-OBJECT-TYPE></SPEC-TYPES>
    <SPEC-OBJECTS>
      <SPEC-OBJECT IDENTIFIER="REQ-1" LONG-NAME="Health"><VALUES>
        <ATTRIBUTE-VALUE-STRING THE-VALUE="returns 200"><DEFINITION>
          <ATTRIBUTE-DEFINITION-STRING-REF>ATTR-1</ATTRIBUTE-DEFINITION-STRING-REF>
        </DEFINITION></ATTRIBUTE-VALUE-STRING>
      </VALUES></SPEC-OBJECT>
      <SPEC-OBJECT IDENTIFIER="REQ-2" LONG-NAME="Database"><VALUES>
        <ATTRIBUTE-VALUE-STRING THE-VALUE="status is ready"><DEFINITION>
          <ATTRIBUTE-DEFINITION-STRING-REF>ATTR-1</ATTRIBUTE-DEFINITION-STRING-REF>
        </DEFINITION></ATTRIBUTE-VALUE-STRING>
      </VALUES></SPEC-OBJECT>
    </SPEC-OBJECTS>
  </REQ-IF-CONTENT></CORE-CONTENT>
</REQ-IF>'''
        result = prepare_request(
            files=[InputFile("requirements.reqif", reqif)],
            target=_target(),
            selections=_selections(),
        )
        self.assertEqual(
            [item.requirement_id for item in result.request.requirements],
            ["REQ-1", "REQ-2"],
        )
        self.assertIn("Description: returns 200", result.request.requirements[0].content)

        malicious = b'<!DOCTYPE root [<!ENTITY xxe SYSTEM "file:///etc/passwd">]><root>&xxe;</root>'
        with self.assertRaises(IngestionError) as raised:
            prepare_request(
                files=[InputFile("malicious.xml", malicious)],
                target=_target(),
                selections=_selections(),
            )
        self.assertEqual(raised.exception.code, "XML_DTD_REJECTED")

        late_dtd = b" " * 70_000 + malicious
        with self.assertRaises(IngestionError) as late_error:
            prepare_request(
                files=[InputFile("late-malicious.xml", late_dtd)],
                target=_target(),
                selections=_selections(),
            )
        self.assertEqual(late_error.exception.code, "XML_DTD_REJECTED")

    def test_yaml_mixed_key_types_do_not_leak_serializer_errors(self):
        yaml_text = """requirements:
  - 1: numeric key
    '1': string key
    content_note: preserved
"""
        result = prepare_request(
            files=[InputFile("mixed-keys.yaml", yaml_text.encode("utf-8"))],
            target=_target(),
            selections=_selections(),
        )
        self.assertIn("numeric key", result.request.requirements[0].content)
        self.assertIn("string key", result.request.requirements[0].content)

    def test_html_removes_executable_content_without_fetching_links(self):
        html = b'''<html><head><style>.hidden{}</style><script>steal()</script></head>
<body><h1>Health</h1><p>Returns 200</p><a href="https://example.invalid">label</a></body></html>'''
        result = prepare_request(
            files=[InputFile("requirements.html", html)],
            target=_target(),
            selections=_selections(),
        )
        content = result.request.requirements[0].content
        self.assertIn("Health", content)
        self.assertIn("Returns 200", content)
        self.assertIn("label", content)
        self.assertNotIn("steal", content)
        self.assertNotIn("example.invalid", content)

    def test_yaml_unsafe_tags_and_office_zip_traversal_are_rejected(self):
        unsafe_yaml = b"!!python/object/apply:os.system ['echo should-not-run']"
        with self.assertRaises(IngestionError) as yaml_error:
            prepare_request(
                files=[InputFile("unsafe.yaml", unsafe_yaml)],
                target=_target(),
                selections=_selections(),
            )
        self.assertEqual(yaml_error.exception.code, "STRUCTURED_TEXT_PARSE_FAILED")

        output = BytesIO()
        with zipfile.ZipFile(output, "w") as archive:
            archive.writestr("../escape", "no")
            archive.writestr("word/document.xml", "<document/>")
        with self.assertRaises(IngestionError) as archive_error:
            prepare_request(
                files=[InputFile("traversal.docx", output.getvalue())],
                target=_target(),
                selections=_selections(),
            )
        self.assertEqual(archive_error.exception.code, "ARCHIVE_PATH_INVALID")

    def test_limits_reject_silent_truncation_and_legacy_binary_office(self):
        rows = ["Summary,Description"] + [f"R{index},value" for index in range(21)]
        with self.assertRaises(IngestionError) as count_error:
            prepare_request(
                files=[InputFile("many.csv", "\n".join(rows).encode("utf-8"))],
                target=_target(),
                selections=_selections(),
            )
        self.assertEqual(count_error.exception.code, "TOO_MANY_REQUIREMENTS")

        with self.assertRaises(IngestionError) as size_error:
            prepare_request(
                frontend_text="abcdef",
                target=_target(),
                selections=_selections(),
                limits=IngestionLimits(max_requirement_bytes=5),
            )
        self.assertEqual(size_error.exception.code, "REQUIREMENT_TOO_LARGE")

        with self.assertRaises(IngestionError) as legacy_error:
            prepare_request(
                files=[InputFile("legacy.doc", b"not a modern Office document")],
                target=_target(),
                selections=_selections(),
            )
        self.assertEqual(legacy_error.exception.code, "LEGACY_OFFICE_UNSUPPORTED")

    def test_workflow_exposes_ingestion_without_changing_generate_contract(self):
        workflow = IntentToExecutionWorkflow(
            SimpleNamespace(),
            SimpleNamespace(compiler=SimpleNamespace()),
            coordinator=SimpleNamespace(),
        )
        result = workflow.prepare_design_request(
            frontend_text="health returns 200",
            target=_target(),
            selections=_selections(),
            request_id="REQ-FRONTEND-1",
        )

        self.assertEqual(result.request.request_id, "REQ-FRONTEND-1")
        self.assertEqual(result.request.requirements[0].content, "health returns 200")


if __name__ == "__main__":
    unittest.main()
